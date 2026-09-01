from __future__ import annotations

import io
from dataclasses import dataclass

import fitz  # PyMuPDF #add fitz dependency
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation


@dataclass
class ExtractedDocument:
    text: str
    metadata: dict


class DocumentExtractor:
    """
    Extract text from supported document formats.

    Supported:
    - PDF
    - DOCX
    - PPTX
    - TXT
    - Images
    """

    async def extract(
        self,
        file_bytes: bytes,
        filename: str,
        mime_type: str,
    ) -> ExtractedDocument:

        extension = filename.lower().rsplit(".", 1)[-1]

        if extension == "pdf":
            return self._extract_pdf(file_bytes)

        if extension == "docx":
            return self._extract_docx(file_bytes)

        if extension == "pptx":
            return self._extract_pptx(file_bytes)

        if extension == "txt":
            return self._extract_txt(file_bytes)

        if mime_type.startswith("image/"):
            return self._extract_image(file_bytes)

        raise ValueError(
            f"Unsupported file type: {filename}"
        )

    def _extract_pdf(
        self,
        file_bytes: bytes,
    ) -> ExtractedDocument:

        document = fitz.open(
            stream=file_bytes,
            filetype="pdf",
        )

        pages = []
        page_numbers = []

        for page_number, page in enumerate(document, start=1):
            text = page.get_text("text").strip()

            if text:
                pages.append(
                    f"[Page {page_number}]\n{text}"
                )
                page_numbers.append(page_number)

        document.close()

        return ExtractedDocument(
            text="\n\n".join(pages),
            metadata={
                "file_type": "pdf",
                "page_count": len(page_numbers),
            },
        )

    def _extract_docx(
        self,
        file_bytes: bytes,
    ) -> ExtractedDocument:

        document = Document(
            io.BytesIO(file_bytes)
        )

        paragraphs = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()

            if text:
                paragraphs.append(text)

        return ExtractedDocument(
            text="\n\n".join(paragraphs),
            metadata={
                "file_type": "docx",
            },
        )

    def _extract_pptx(
        self,
        file_bytes: bytes,
    ) -> ExtractedDocument:

        presentation = Presentation(
            io.BytesIO(file_bytes)
        )

        slides = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):

            texts = []

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue

                text = shape.text.strip()

                if text:
                    texts.append(text)

            if texts:
                slides.append(
                    f"[Slide {slide_number}]\n"
                    + "\n".join(texts)
                )

        return ExtractedDocument(
            text="\n\n".join(slides),
            metadata={
                "file_type": "pptx",
                "slide_count": len(
                    presentation.slides
                ),
            },
        )

    @staticmethod
    def _extract_txt(
        file_bytes: bytes,
    ) -> ExtractedDocument:

        text = file_bytes.decode(
            "utf-8",
            errors="ignore",
        )

        return ExtractedDocument(
            text=text,
            metadata={
                "file_type": "txt",
            },
        )

    @staticmethod
    def _extract_image(
        file_bytes: bytes,
    ) -> ExtractedDocument:

        image = Image.open(
            io.BytesIO(file_bytes)
        )

        text = pytesseract.image_to_string(image)

        return ExtractedDocument(
            text=text,
            metadata={
                "file_type": "image",
                "ocr": True,
            },
        )
